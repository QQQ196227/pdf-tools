import sys
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import io

def pdf_to_pptx(input_path, output_path):
    pdf_doc = fitz.open(input_path)
    prs = Presentation()

    # 设置幻灯片尺寸为 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for page_num in range(len(pdf_doc)):
        page = pdf_doc[page_num]

        # 将 PDF 页面渲染为高清图片（2倍分辨率）
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat)
        img_bytes = pix.tobytes("png")

        # 添加空白幻灯片
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 计算图片尺寸，保持宽高比
        img_width = pix.width
        img_height = pix.height
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # 缩放图片以适应幻灯片
        ratio = min(slide_w / img_width, slide_h / img_height)
        display_w = int(img_width * ratio)
        display_h = int(img_height * ratio)

        # 居中放置
        left = int((slide_w - display_w) / 2)
        top = int((slide_h - display_h) / 2)

        # 将图片嵌入幻灯片
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, left, top, display_w, display_h)

    pdf_doc.close()
    prs.save(output_path)
    return True

if __name__ == '__main__':
    if len(sys.argv) != 3:
        print("Usage: pdf_to_pptx.py input.pdf output.pptx")
        sys.exit(1)
    try:
        pdf_to_pptx(sys.argv[1], sys.argv[2])
        print("OK")
    except Exception as e:
        print(f"ERROR: {e}", file=sys.stderr)
        sys.exit(1)
